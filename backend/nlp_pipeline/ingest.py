"""Этап 1 пайплайна: сырой файл (docx/pptx/pdf) -> список текстовых блоков + метаданные.

Для PDF со сканами (типично для старых отчётов в горно-металлургической отрасли)
используется OCR-фолбэк через tesseract (rus+eng), если текстовый слой страницы
пуст или почти пуст. Метаданные файла (имя, дата изменения) — заготовка под
source/date в модели верификации фактов (schema.py: Entity.attrs).
"""
from __future__ import annotations

import datetime
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


@dataclass
class TextBlock:
    text: str
    kind: str  # "paragraph" | "table" | "heading"
    location: str  # "para 12" / "slide 3" / "page 7"
    meta: dict[str, Any] = field(default_factory=dict)


@dataclass
class FileMeta:
    source_file: str
    path: str
    modified: str
    ext: str


def file_meta(path: str | Path) -> FileMeta:
    p = Path(path)
    stat = p.stat()
    return FileMeta(
        source_file=p.name,
        path=str(p),
        modified=datetime.datetime.fromtimestamp(stat.st_mtime).date().isoformat(),
        ext=p.suffix.lower(),
    )


def load_docx(path: str | Path) -> list[TextBlock]:
    import docx

    doc = docx.Document(str(path))
    blocks: list[TextBlock] = []
    for i, p in enumerate(doc.paragraphs):
        text = p.text.strip()
        if not text:
            continue
        style = (p.style.name or "").lower() if p.style else ""
        kind = "heading" if "heading" in style or "title" in style else "paragraph"
        blocks.append(TextBlock(text=text, kind=kind, location=f"para {i}"))
    for ti, table in enumerate(doc.tables):
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            blocks.append(TextBlock(text="\n".join(rows), kind="table", location=f"table {ti}"))
    return blocks


def load_pptx(path: str | Path) -> list[TextBlock]:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[TextBlock] = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False):
                rows = []
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    texts.append("[таблица]\n" + "\n".join(rows))
        if texts:
            blocks.append(TextBlock(text="\n".join(texts), kind="paragraph", location=f"slide {i}"))
    return blocks


def load_pdf(path: str | Path, ocr_lang: str = "rus+eng", min_chars_for_text_layer: int = 20) -> list[TextBlock]:
    """Текстовый слой через PyMuPDF; если страница почти без текста (скан) —
    рендерим в изображение и распознаём через tesseract."""
    import fitz  # PyMuPDF

    blocks: list[TextBlock] = []
    doc = fitz.open(str(path))
    try:
        for i, page in enumerate(doc):
            text = (page.get_text("text") or "").strip()
            used_ocr = False
            if len(text) < min_chars_for_text_layer:
                ocr_text = _ocr_page(page, ocr_lang)
                if ocr_text:
                    text = ocr_text
                    used_ocr = True
            if text:
                blocks.append(TextBlock(text=text, kind="paragraph", location=f"page {i + 1}", meta={"ocr": used_ocr}))
    finally:
        doc.close()
    return blocks


def _ocr_page(page, ocr_lang: str) -> str:
    try:
        import pytesseract
        from PIL import Image

        pix = page.get_pixmap(dpi=200)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        return pytesseract.image_to_string(img, lang=ocr_lang).strip()
    except Exception as e:
        print(f"[ingest] OCR не удался для страницы: {e}")
        return ""


LOADERS = {".docx": load_docx, ".pptx": load_pptx, ".pdf": load_pdf}


def load_document(path: str | Path) -> tuple[list[TextBlock], FileMeta]:
    meta = file_meta(path)
    loader = LOADERS.get(meta.ext)
    if not loader:
        raise ValueError(f"Неподдерживаемый формат документа: {meta.ext}, {path}")
    return loader(path), meta
