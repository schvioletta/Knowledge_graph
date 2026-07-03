"""Отбор ограниченного числа файлов на категорию — чтобы строить граф не по
всему корпусу, а по репрезентативной выборке (например, 3 файла из каждой
категории документов), когда корпус большой и гонять LLM по всем файлам
дорого.

Категория определяется по имени подпапки внутри корня сырых данных
(`data/raw/статьи/...`, `data/raw/доклады/...`) — так пользователь явно
контролирует деление на категории, без хрупкой автоклассификации по
содержимому. Если файлы лежат плоско (без подпапок, как в этом репозитории
по умолчанию) — категория определяется по расширению как разумный фолбэк.
"""
from __future__ import annotations

import random
from collections import defaultdict
from pathlib import Path

CATEGORY_BY_EXT = {
    ".pptx": "доклад",
    ".docx": "статья_или_отчёт",
    ".pdf": "статья_или_отчёт_pdf",
}

STRATEGIES = ("random", "largest", "newest", "most_tables")


def infer_category(path: Path, raw_root: Path) -> str:
    """Категория — путь подпапок от raw_root, без самой глубокой (обычно это
    год/выпуск, а не тематическая категория). Иначе на корпусах вида
    `Журналы/Цветные металлы/2005-2/...` каждый год/выпуск считался бы
    отдельной категорией, и --per-category N раздувался бы до N файлов на
    каждый такой год вместо N файлов на «Цветные металлы» в целом."""
    parent = path.resolve().parent
    root = raw_root.resolve()
    if parent == root:
        return CATEGORY_BY_EXT.get(path.suffix.lower(), "прочее")
    rel_parts = parent.relative_to(root).parts
    if len(rel_parts) > 1:
        rel_parts = rel_parts[:-1]
    return "/".join(rel_parts)


def _count_tables_cheap(path: Path) -> int:
    """Дешёвая оценка «информационной плотности» без полного парсинга — просто
    считает конверты таблиц/файл, чтобы не тратить время на файлах, которые
    всё равно не попадут в выборку."""
    try:
        if path.suffix.lower() == ".docx":
            import docx
            return len(docx.Document(str(path)).tables)
        if path.suffix.lower() == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(path))
            return sum(1 for slide in prs.slides for shape in slide.shapes if getattr(shape, "has_table", False))
    except Exception:
        pass
    return 0


def select_files(
    files: list[str | Path],
    raw_root: str | Path,
    per_category: int = 3,
    strategy: str = "random",
    seed: int = 42,
) -> dict[str, list[Path]]:
    if strategy not in STRATEGIES:
        raise ValueError(f"Неизвестная стратегия отбора: {strategy!r}, допустимо {STRATEGIES}")

    raw_root = Path(raw_root)
    groups: dict[str, list[Path]] = defaultdict(list)
    for f in files:
        p = Path(f)
        if not p.is_file():
            continue
        groups[infer_category(p, raw_root)].append(p)

    rng = random.Random(seed)
    selected: dict[str, list[Path]] = {}
    for category, flist in groups.items():
        if strategy == "random":
            chosen = rng.sample(flist, min(per_category, len(flist)))
        elif strategy == "largest":
            chosen = sorted(flist, key=lambda p: p.stat().st_size, reverse=True)[:per_category]
        elif strategy == "newest":
            chosen = sorted(flist, key=lambda p: p.stat().st_mtime, reverse=True)[:per_category]
        else:  # most_tables
            chosen = sorted(flist, key=_count_tables_cheap, reverse=True)[:per_category]
        selected[category] = chosen
    return selected
